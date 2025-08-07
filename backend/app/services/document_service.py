import uuid
from datetime import datetime
from typing import Dict, Any, Optional, BinaryIO
import logging

# Document processing libraries
import PyPDF2
from docx import Document
import openpyxl
from pptx import Presentation
import zipfile
import xml.etree.ElementTree as ET

logger = logging.getLogger(__name__)

class DocumentService:
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
    
    TEXT_FILE_TYPES = [
        'text/plain', 'text/csv', 'application/json', 'text/markdown',
        'text/html', 'application/javascript', 'text/css', 'application/xml',
        'text/xml', 'application/yaml', 'text/yaml'
    ]
    
    BINARY_FILE_TYPES = [
        'application/pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document',  # .docx
        'application/vnd.openxmlformats-officedocument.presentationml.presentation',  # .pptx
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',  # .xlsx
        'application/vnd.ms-excel'  # .xls
    ]
    
    def __init__(self):
        logger.info("DocumentService initialized")
    
    async def process_file(self, file_data: bytes, filename: str, content_type: str) -> Dict[str, Any]:
        """Process uploaded file and extract content"""
        logger.info(f"🔍 Processing file: {filename} ({len(file_data)} bytes)")
        
        # Validate file size
        if len(file_data) > self.MAX_FILE_SIZE:
            raise ValueError(f"File size exceeds {self.MAX_FILE_SIZE / 1024 / 1024}MB limit")
        
        # Detect actual file type if needed
        if not content_type or content_type == 'application/octet-stream':
            content_type = self._detect_content_type(file_data, filename)
        
        # Validate file type
        if not self._is_supported_type(content_type, filename):
            extension = filename.lower().split('.')[-1] if '.' in filename else 'unknown'
            supported_exts = ', '.join(self._get_supported_extensions())
            raise ValueError(f"Unsupported file type: .{extension}\\n\\nSupported formats: {supported_exts}")
        
        try:
            content = ""
            extracted_text = ""
            metadata = {}
            
            if self._is_text_file(content_type, filename):
                content = await self._extract_text_content(file_data)
                extracted_text = content
                metadata = {
                    'type': 'Text File',
                    'extraction_method': 'Direct text reading',
                    'extraction_success': True,
                    'extraction_quality': 'excellent'
                }
            elif self._is_binary_file(content_type, filename):
                result = await self._extract_binary_content(file_data, filename)
                content = result['text']
                extracted_text = result['text']
                metadata = result['metadata']
            
            summary = self._generate_summary(content, filename, metadata)
            
            document_info = {
                'id': self._generate_id(),
                'name': filename,
                'type': content_type or self._get_type_from_extension(filename),
                'size': len(file_data),
                'uploaded_at': datetime.now(),
                'content': content,
                'extracted_text': extracted_text,
                'summary': summary,
                'metadata': metadata
            }
            
            logger.info(f"✅ Successfully processed {filename}")
            return document_info
            
        except Exception as error:
            logger.error(f"❌ Failed to process {filename}: {error}")
            raise ValueError(f"Failed to process {filename}: {str(error)}")
    
    async def _extract_text_content(self, file_data: bytes) -> str:
        """Extract text from text files"""
        try:
            # Try UTF-8 first, fallback to other encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    return file_data.decode(encoding)
                except UnicodeDecodeError:
                    continue
                    
            # If all encodings fail, decode with errors='replace'
            return file_data.decode('utf-8', errors='replace')
            
        except Exception as error:
            raise ValueError(f"Failed to read text file: {error}")
    
    async def _extract_binary_content(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """Extract content from binary files"""
        extension = filename.lower().split('.')[-1] if '.' in filename else ''
        
        logger.info(f"📄 Extracting binary content from .{extension} file")
        
        try:
            if extension == 'pdf':
                return await self._extract_pdf_content(file_data, filename)
            elif extension == 'docx':
                return await self._extract_word_content(file_data, filename)
            elif extension == 'pptx':
                return await self._extract_powerpoint_content(file_data, filename)
            elif extension in ['xlsx', 'xls']:
                return await self._extract_excel_content(file_data, filename)
            else:
                raise ValueError(f"Unsupported binary format: .{extension}")
                
        except Exception as error:
            logger.error(f"❌ Error extracting content from {filename}: {error}")
            raise error
    
    async def _extract_pdf_content(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """Extract text from PDF files"""
        logger.info(f"🔍 Starting PDF extraction for: {filename}")
        
        try:
            from io import BytesIO
            pdf_file = BytesIO(file_data)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            num_pages = len(pdf_reader.pages)
            logger.info(f"📖 PDF loaded: {num_pages} pages")
            
            full_text = ""
            total_chars = 0
            pages_with_text = 0
            
            # Process each page
            for page_num in range(num_pages):
                try:
                    logger.debug(f"📄 Processing page {page_num + 1}...")
                    
                    page = pdf_reader.pages[page_num]
                    page_text = page.extract_text()
                    
                    if page_text.strip():
                        full_text += f"\\n--- PAGE {page_num + 1} ---\\n{page_text.strip()}\\n"
                        total_chars += len(page_text)
                        pages_with_text += 1
                        logger.debug(f"✅ Page {page_num + 1}: {len(page_text)} characters")
                    else:
                        full_text += f"\\n--- PAGE {page_num + 1} ---\\n[No text content found]\\n"
                        logger.debug(f"⚠️ Page {page_num + 1}: No text found")
                        
                except Exception as page_error:
                    logger.error(f"❌ Error on page {page_num + 1}: {page_error}")
                    full_text += f"\\n--- PAGE {page_num + 1} ---\\n[Error extracting page: {page_error}]\\n"
            
            # Add summary
            full_text += f"\\n\\n=== PDF SUMMARY ===\\n"
            full_text += f"File: {filename}\\n"
            full_text += f"Total Pages: {num_pages}\\n"
            full_text += f"Pages with Text: {pages_with_text}\\n"
            full_text += f"Total Characters: {total_chars}\\n"
            
            # Determine quality
            quality = 'poor'
            if pages_with_text > 0 and total_chars > 100:
                text_ratio = pages_with_text / num_pages
                if text_ratio >= 0.8 and total_chars > 1000:
                    quality = 'excellent'
                elif text_ratio >= 0.5 and total_chars > 500:
                    quality = 'good'
                elif text_ratio >= 0.2 and total_chars > 100:
                    quality = 'partial'
            
            metadata = {
                'type': 'PDF Document',
                'extraction_method': 'PyPDF2 Text Extraction',
                'page_count': num_pages,
                'pages_with_text': pages_with_text,
                'total_characters': total_chars,
                'extraction_success': pages_with_text > 0,
                'extraction_quality': quality,
                'text_coverage': round((pages_with_text / num_pages) * 100)
            }
            
            logger.info(f"🎯 PDF extraction complete: {quality} quality, {pages_with_text}/{num_pages} pages with text")
            
            if pages_with_text == 0:
                raise ValueError('No text found in PDF. This may be a scanned document or image-based PDF.')
            
            return {'text': full_text, 'metadata': metadata}
            
        except Exception as error:
            logger.error(f"❌ PDF extraction failed: {error}")
            
            # Return meaningful error document
            error_text = f"""PDF Document: {filename}

❌ Extraction Failed: {error}

This PDF could not be processed. Common reasons:
• Scanned document (requires OCR)
• Password protected
• Corrupted file
• Complex formatting

File size: {len(file_data) // 1024}KB
Upload date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

The AI can still reference this document by name but won't have access to its content."""
            
            return {
                'text': error_text,
                'metadata': {
                    'type': 'PDF Document',
                    'extraction_method': 'Failed',
                    'extraction_success': False,
                    'extraction_quality': 'poor',
                    'error': str(error)
                }
            }
    
    async def _extract_word_content(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """Extract text from Word documents"""
        logger.info(f"📝 Starting Word extraction for: {filename}")
        
        try:
            from io import BytesIO
            doc_file = BytesIO(file_data)
            doc = Document(doc_file)
            
            # Extract text from all paragraphs
            text_content = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text.strip())
            
            full_text = '\\n\\n'.join(text_content)
            word_count = len(full_text.split()) if full_text else 0
            paragraph_count = len(text_content)
            
            metadata = {
                'type': 'Word Document (.docx)',
                'extraction_method': 'python-docx',
                'word_count': word_count,
                'paragraphs': paragraph_count,
                'extraction_success': True,
                'extraction_quality': 'excellent',
                'character_count': len(full_text)
            }
            
            logger.info(f"✅ Word extraction complete: {word_count} words")
            
            return {'text': full_text, 'metadata': metadata}
            
        except Exception as error:
            logger.error(f"❌ Word extraction error: {error}")
            raise ValueError(f"Word document extraction failed: {error}")
    
    async def _extract_powerpoint_content(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """Extract text from PowerPoint presentations"""
        logger.info(f"🎯 Starting PowerPoint extraction for: {filename}")
        
        try:
            from io import BytesIO
            ppt_file = BytesIO(file_data)
            presentation = Presentation(ppt_file)
            
            full_text = f"PowerPoint Presentation: {filename}\\n{'=' * 50}\\n\\n"
            slide_count = len(presentation.slides)
            total_text_length = 0
            
            for i, slide in enumerate(presentation.slides, 1):
                slide_text = f"--- SLIDE {i} ---\\n"
                slide_content = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                if slide_content:
                    slide_text += '\\n'.join(slide_content) + '\\n\\n'
                    total_text_length += len('\\n'.join(slide_content))
                else:
                    slide_text += '[No text content]\\n\\n'
                
                full_text += slide_text
            
            metadata = {
                'type': 'PowerPoint Presentation (.pptx)',
                'extraction_method': 'python-pptx',
                'slide_count': slide_count,
                'total_text_length': total_text_length,
                'extraction_success': total_text_length > 0,
                'extraction_quality': 'excellent' if total_text_length > 500 else 'good' if total_text_length > 100 else 'partial'
            }
            
            logger.info(f"✅ PowerPoint extraction complete: {slide_count} slides")
            
            return {'text': full_text, 'metadata': metadata}
            
        except Exception as error:
            logger.error(f"❌ PowerPoint extraction error: {error}")
            raise ValueError(f"PowerPoint extraction failed: {error}")
    
    async def _extract_excel_content(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """Extract content from Excel files"""
        logger.info(f"📊 Starting Excel extraction for: {filename}")
        
        try:
            from io import BytesIO
            excel_file = BytesIO(file_data)
            workbook = openpyxl.load_workbook(excel_file)
            
            full_text = f"Excel Spreadsheet: {filename}\\n{'=' * 50}\\n\\n"
            total_rows = 0
            sheet_count = len(workbook.worksheets)
            
            for worksheet in workbook.worksheets:
                sheet_name = worksheet.title
                full_text += f"--- SHEET: {sheet_name} ---\\n"
                
                rows_in_sheet = 0
                for row in worksheet.iter_rows(values_only=True):
                    # Filter out empty rows
                    row_data = [str(cell) if cell is not None else '' for cell in row]
                    if any(cell.strip() for cell in row_data):
                        full_text += f"Row {rows_in_sheet + 1}: {' | '.join(row_data)}\\n"
                        rows_in_sheet += 1
                        total_rows += 1
                
                full_text += '\\n'
            
            metadata = {
                'type': 'Excel Spreadsheet',
                'extraction_method': 'openpyxl',
                'sheet_count': sheet_count,
                'total_rows': total_rows,
                'extraction_success': True,
                'extraction_quality': 'excellent'
            }
            
            logger.info(f"✅ Excel extraction complete: {sheet_count} sheets, {total_rows} rows")
            
            return {'text': full_text, 'metadata': metadata}
            
        except Exception as error:
            logger.error(f"❌ Excel extraction error: {error}")
            raise ValueError(f"Excel extraction failed: {error}")
    
    def _detect_content_type(self, file_data: bytes, filename: str) -> str:
        """Detect file content type from extension (simplified version)"""
        # For now, use extension-based detection
        # In production, you might want to install libmagic for better detection
        return self._get_type_from_extension(filename)
    
    def _is_text_file(self, mime_type: str, filename: str) -> bool:
        """Check if file is a text file"""
        if mime_type in self.TEXT_FILE_TYPES:
            return True
        
        extension = filename.lower().split('.')[-1] if '.' in filename else ''
        text_extensions = [
            'txt', 'md', 'json', 'csv', 'html', 'htm', 'js', 'ts', 'jsx', 'tsx',
            'css', 'scss', 'sass', 'xml', 'yaml', 'yml', 'log', 'sql', 'py', 'java',
            'cpp', 'c', 'h', 'php', 'rb', 'go', 'rs', 'sh', 'bat', 'ps1', 'r'
        ]
        
        return extension in text_extensions
    
    def _is_binary_file(self, mime_type: str, filename: str) -> bool:
        """Check if file is a supported binary file"""
        if mime_type in self.BINARY_FILE_TYPES:
            return True
        
        extension = filename.lower().split('.')[-1] if '.' in filename else ''
        binary_extensions = ['pdf', 'docx', 'pptx', 'xlsx', 'xls']
        
        return extension in binary_extensions
    
    def _is_supported_type(self, mime_type: str, filename: str) -> bool:
        """Check if file type is supported"""
        return self._is_text_file(mime_type, filename) or self._is_binary_file(mime_type, filename)
    
    def _get_type_from_extension(self, filename: str) -> str:
        """Get MIME type from file extension"""
        extension = filename.lower().split('.')[-1] if '.' in filename else ''
        
        type_map = {
            'txt': 'text/plain',
            'md': 'text/markdown',
            'json': 'application/json',
            'csv': 'text/csv',
            'html': 'text/html',
            'htm': 'text/html',
            'js': 'application/javascript',
            'ts': 'application/typescript',
            'css': 'text/css',
            'xml': 'application/xml',
            'yaml': 'application/yaml',
            'yml': 'application/yaml',
            'pdf': 'application/pdf',
            'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'xls': 'application/vnd.ms-excel'
        }
        
        return type_map.get(extension, 'application/octet-stream')
    
    def _generate_summary(self, content: str, filename: str, metadata: Dict[str, Any] = None) -> str:
        """Generate a summary of the document"""
        if metadata is None:
            metadata = {}
        
        lines = content.split('\\n') if content else []
        word_count = len(content.split()) if content else 0
        char_count = len(content)
        
        summary = f"📁 File: {filename}\\n"
        
        if metadata.get('type'):
            summary += f"📄 Type: {metadata['type']}\\n"
        
        summary += f"📊 Size: {char_count:,} characters, {word_count:,} words, {len(lines):,} lines\\n"
        
        # Add specific metadata
        if metadata.get('page_count'):
            summary += f"📄 Pages: {metadata['page_count']}"
            if metadata.get('pages_with_text'):
                summary += f" ({metadata['pages_with_text']} with text, {metadata.get('text_coverage', 0)}% coverage)"
            summary += '\\n'
        
        if metadata.get('slide_count'):
            summary += f"🎯 Slides: {metadata['slide_count']}\\n"
        
        if metadata.get('sheet_count'):
            summary += f"📊 Sheets: {metadata['sheet_count']}\\n"
        
        # Quality indicator
        if metadata.get('extraction_quality'):
            quality_emoji = {
                'excellent': '🟢',
                'good': '🟡',
                'partial': '🟠',
                'poor': '🔴'
            }
            quality = metadata['extraction_quality']
            summary += f"{quality_emoji.get(quality, '⚪')} Quality: {quality}\\n"
        
        # Status
        if metadata.get('extraction_success') is False:
            summary += "⚠️ Extraction failed\\n"
        else:
            summary += "✅ Content extracted successfully\\n"
        
        # Preview
        if content and len(content) > 100 and metadata.get('extraction_success') is not False:
            preview = content[:200].strip()
            summary += f"\\n📖 Preview: {preview}{'...' if len(content) > 200 else ''}"
        
        return summary
    
    def _generate_id(self) -> str:
        """Generate a unique document ID"""
        return str(uuid.uuid4())
    
    def format_document_for_ai(self, document: Dict[str, Any]) -> str:
        """Format document for AI consumption"""
        uploaded_at = document.get('uploaded_at')
        if isinstance(uploaded_at, datetime):
            uploaded_at_str = uploaded_at.strftime('%Y-%m-%d %H:%M:%S')
        else:
            uploaded_at_str = str(uploaded_at)
        
        formatted_doc = f"""📄 DOCUMENT: {document.get('name', 'Unknown')}
📋 Type: {document.get('type', 'Unknown')}
📅 Uploaded: {uploaded_at_str}
📊 Summary: {document.get('summary', 'No summary available')}"""
        
        content = document.get('extracted_text') or document.get('content', '')
        formatted_doc += f"\\n\\n📖 CONTENT:\\n{'-' * 80}\\n{content}"
        formatted_doc += f"\\n{'-' * 80}\\n📄 END OF DOCUMENT: {document.get('name', 'Unknown')}"
        
        return formatted_doc
    
    def _get_supported_extensions(self) -> list:
        """Get list of supported file extensions"""
        return [
            # Text files
            'txt', 'md', 'json', 'csv', 'html', 'htm', 'js', 'ts', 'jsx', 'tsx',
            'css', 'scss', 'sass', 'xml', 'yaml', 'yml', 'log', 'sql', 'py', 'java',
            'cpp', 'c', 'h', 'php', 'rb', 'go', 'rs', 'sh', 'bat', 'ps1', 'r',
            # Binary files
            'pdf', 'docx', 'pptx', 'xlsx', 'xls'
        ]
    
    async def delete_document(self, document_id: str, user_id: str) -> bool:
        """Delete a document from the database (user must own the document)"""
        logger.info(f"🗑️ Deleting document {document_id} for user {user_id}")
        
        try:
            # Import here to avoid circular import
            from supabase import create_client
            from ..core.config import settings
            
            # Initialize Supabase client
            if not settings.SUPABASE_SERVICE_ROLE_KEY:
                raise ValueError("Supabase service role key not configured")
            supabase = create_client(settings.SUPABASE_URL, settings.SUPABASE_SERVICE_ROLE_KEY)
            
            # First, verify the document belongs to the user's agent
            document_check = supabase.table("agent_documents").select("id, agent_id").eq("id", document_id).execute()
            
            if not document_check.data:
                raise ValueError(f"Document {document_id} not found")
            
            agent_id = document_check.data[0]["agent_id"]
            
            # Verify the agent belongs to the user
            agent_check = supabase.table("user_agents").select("id").eq("id", agent_id).eq("user_id", user_id).execute()
            
            if not agent_check.data:
                raise ValueError(f"Document {document_id} does not belong to user {user_id}")
            
            # Delete the document
            delete_result = supabase.table("agent_documents").delete().eq("id", document_id).execute()
            
            if not delete_result.data:
                raise ValueError(f"Failed to delete document {document_id}")
            
            logger.info(f"✅ Successfully deleted document {document_id}")
            return True
            
        except Exception as e:
            logger.error(f"❌ Error deleting document {document_id}: {e}")
            raise

# Create singleton instance
document_service = DocumentService()